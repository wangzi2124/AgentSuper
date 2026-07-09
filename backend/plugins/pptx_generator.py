"""
PPTX Generator Plugin

Creates PowerPoint presentations (.pptx) from structured content using python-pptx.
"""
import json
import os
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

PLUGIN_NAME = "pptx-generator"
PLUGIN_VERSION = "0.1.0"
PLUGIN_DESCRIPTION = "Creates PowerPoint presentations (.pptx) from structured content"


def _hex_to_rgb(hex_str: str) -> RGBColor:
    h = hex_str.lstrip("#")
    return RGBColor(int(h[0:2], 16), int(h[2:4], 16), int(h[4:6], 16))


def _set_slide_bg(slide, color_hex: str):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = _hex_to_rgb(color_hex)


def _add_textbox(slide, left, top, width, height, text, font_size=14,
                 bold=False, color="000000", alignment=PP_ALIGN.LEFT,
                 font_name="Calibri", italic=False):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = _hex_to_rgb(color)
    p.font.name = font_name
    p.font.italic = italic
    p.alignment = alignment
    return txBox


def _add_bullet_textbox(slide, left, top, width, height, items, font_size=14,
                        color="000000", font_name="Calibri"):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.text = item
        p.font.size = Pt(font_size)
        p.font.color.rgb = _hex_to_rgb(color)
        p.font.name = font_name
        p.level = 0
        p.space_after = Pt(4)
    return txBox


def _add_table(slide, left, top, width, height, headers, rows,
               header_bg="4472C4", header_color="FFFFFF", alt_color="D9E2F3"):
    col_count = len(headers)
    row_count = 1 + len(rows)
    table_shape = slide.shapes.add_table(row_count, col_count,
                                         Inches(left), Inches(top),
                                         Inches(width), Inches(height))
    table = table_shape.table

    for ci, h in enumerate(headers):
        cell = table.cell(0, ci)
        cell.text = str(h)
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(11)
            p.font.bold = True
            p.font.color.rgb = _hex_to_rgb(header_color)
            p.font.name = "Calibri"
            p.alignment = PP_ALIGN.CENTER
        cell.fill.solid()
        cell.fill.fore_color.rgb = _hex_to_rgb(header_bg)
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE

    for ri, row in enumerate(rows):
        for ci, val in enumerate(row):
            cell = table.cell(ri + 1, ci)
            cell.text = str(val)
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(10)
                p.font.name = "Calibri"
                p.alignment = PP_ALIGN.CENTER
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            if ri % 2 == 1:
                cell.fill.solid()
                cell.fill.fore_color.rgb = _hex_to_rgb(alt_color)

    return table_shape


def tool_create_pptx(title: str = "Presentation", slides: str = "[]",
                     output_path: str = "") -> str:
    """
    Create a PowerPoint presentation (.pptx) with structured slides.

    Parameters:
    - title: presentation title (used for filename if output_path not given)
    - slides: JSON array of slide objects. Each object has:
        {"type": "title", "text": "...", "subtitle": "...",
         "bg_color": "optional hex", "font_color": "optional hex"}
        {"type": "content", "title": "...", "bullets": ["...", "..."],
         "bg_color": "optional hex", "font_color": "optional hex"}
        {"type": "two_column", "title": "...",
         "left_bullets": ["..."], "right_bullets": ["..."],
         "bg_color": "optional hex"}
        {"type": "table", "title": "...",
         "headers": [...], "rows": [[...], ...],
         "bg_color": "optional hex"}
        {"type": "section_header", "text": "...",
         "bg_color": "optional hex", "font_color": "optional hex"}
    - output_path: optional absolute path to save (auto-generated if empty)
    """
    base_dir = Path(__file__).resolve().parents[1]
    output_dir = base_dir / "data" / "generated"
    output_dir.mkdir(parents=True, exist_ok=True)

    if output_path:
        p = Path(output_path)
        if not p.is_absolute():
            output_path = str(output_dir / p)
    else:
        from datetime import datetime
        safe_title = "".join(c for c in title if c not in '<>:"/\\|?*').strip()
        if not safe_title:
            safe_title = "presentation"
        ts = datetime.now().strftime("%Y%m%d_%H%M%S")
        output_path = str(output_dir / f"{safe_title}_{ts}")

    if not output_path.lower().endswith(".pptx"):
        output_path += ".pptx"

    try:
        parsed_slides = json.loads(slides)
    except json.JSONDecodeError as e:
        return f"Error: invalid slides JSON - {e}"

    if not parsed_slides:
        return "Error: at least one slide is required"

    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    blank_layout = prs.slide_layouts[6]

    for slide_data in parsed_slides:
        slide_type = slide_data.get("type", "content")
        bg_color = slide_data.get("bg_color", "")
        font_color = slide_data.get("font_color", "000000")

        slide = prs.slides.add_slide(blank_layout)

        if bg_color:
            _set_slide_bg(slide, bg_color)

        if slide_type == "title":
            text = slide_data.get("text", title)
            subtitle = slide_data.get("subtitle", "")
            fc = slide_data.get("font_color", "FFFFFF") if bg_color else "000000"
            if subtitle:
                _add_textbox(slide, 1.5, 2.0, 10.3, 1.5, text,
                             font_size=44, bold=True, color=fc,
                             alignment=PP_ALIGN.CENTER)
                _add_textbox(slide, 1.5, 3.8, 10.3, 1.0, subtitle,
                             font_size=20, color=fc,
                             alignment=PP_ALIGN.CENTER)
            else:
                _add_textbox(slide, 1.5, 2.5, 10.3, 1.5, text,
                             font_size=44, bold=True, color=fc,
                             alignment=PP_ALIGN.CENTER)

        elif slide_type == "section_header":
            text = slide_data.get("text", "")
            fc = slide_data.get("font_color", "FFFFFF") if bg_color else "000000"
            _add_textbox(slide, 1.5, 2.5, 10.3, 1.5, text,
                         font_size=40, bold=True, color=fc,
                         alignment=PP_ALIGN.CENTER)

        elif slide_type == "content":
            title_text = slide_data.get("title", "")
            bullets = slide_data.get("bullets", [])
            fc = slide_data.get("font_color", "333333") if not bg_color else \
                slide_data.get("font_color", "FFFFFF")
            if title_text:
                _add_textbox(slide, 0.8, 0.5, 11.7, 0.8, title_text,
                             font_size=28, bold=True, color=fc)
            if bullets:
                _add_bullet_textbox(slide, 0.8, 1.5, 11.7, 5.5, bullets,
                                    font_size=18, color=fc)

        elif slide_type == "two_column":
            title_text = slide_data.get("title", "")
            left_bullets = slide_data.get("left_bullets", [])
            right_bullets = slide_data.get("right_bullets", [])
            fc = slide_data.get("font_color", "333333") if not bg_color else \
                slide_data.get("font_color", "FFFFFF")
            if title_text:
                _add_textbox(slide, 0.8, 0.3, 11.7, 0.7, title_text,
                             font_size=28, bold=True, color=fc)
            if left_bullets:
                _add_textbox(slide, 0.8, 1.3, 5.5, 0.4,
                             slide_data.get("left_header", ""),
                             font_size=20, bold=True, color=fc)
                _add_bullet_textbox(slide, 0.8, 1.8, 5.5, 5.0, left_bullets,
                                    font_size=16, color=fc)
            if right_bullets:
                _add_textbox(slide, 7.0, 1.3, 5.5, 0.4,
                             slide_data.get("right_header", ""),
                             font_size=20, bold=True, color=fc)
                _add_bullet_textbox(slide, 7.0, 1.8, 5.5, 5.0, right_bullets,
                                    font_size=16, color=fc)

        elif slide_type == "table":
            title_text = slide_data.get("title", "")
            headers = slide_data.get("headers", [])
            rows = slide_data.get("rows", [])
            fc = slide_data.get("font_color", "333333") if not bg_color else \
                slide_data.get("font_color", "FFFFFF")
            if title_text:
                _add_textbox(slide, 0.8, 0.3, 11.7, 0.7, title_text,
                             font_size=28, bold=True, color=fc)
            if headers:
                _add_table(slide, 0.8, 1.3, 11.7, 5.5, headers, rows)

    prs.save(output_path)
    return f"Presentation created successfully: {output_path}"
